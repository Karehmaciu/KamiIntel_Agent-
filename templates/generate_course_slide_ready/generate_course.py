
from flask import request, render_template, redirect, flash
from datetime import datetime
from docx import Document
from pptx import Presentation
import os

@app.route("/generate_course", methods=["GET", "POST"])
def generate_course():
    if request.method == "POST":
        prompt = request.form.get("course_prompt", "")
        category = request.form.get("category")

        if not prompt:
            flash("⚠️ Please enter a course topic.", "danger")
            return redirect(request.url)

        # Main course generation
        result = generate_course_outline(prompt)
        if isinstance(result, str):
            flash(result, "danger")
            return redirect(request.url)

        course_outline = result["content"]
        course_category = category or result["category"]
        date_str = datetime.today().strftime('%Y-%m-%d')
        safe_topic = prompt.lower().replace(" ", "_").replace("/", "-")
        base_filename = f"{safe_topic}_{date_str}"

        # Define paths
        save_dir = "static/generated_courses"
        os.makedirs(save_dir, exist_ok=True)
        txt_file = f"{base_filename}.txt"
        docx_file = f"{base_filename}.docx"
        pptx_file = f"{base_filename}.pptx"
        pdf_file = f"{base_filename}.pdf"  # Placeholder

        txt_path = os.path.join(save_dir, txt_file)
        docx_path = os.path.join(save_dir, docx_file)
        pptx_path = os.path.join(save_dir, pptx_file)

        # Save .txt
        with open(txt_path, "w", encoding="utf-8") as f:
            f.write(f"{prompt}
Category: {course_category}
Generated on: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}

{course_outline}")

        # Save .docx
        doc = Document()
        doc.add_heading(prompt, 0)
        doc.add_paragraph(f"Category: {course_category}")
        doc.add_paragraph(f"Generated on: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        doc.add_paragraph("")
        for section in course_outline.strip().split("\n\n"):
            lines = section.strip().split("\n")
            if lines:
                doc.add_heading(lines[0], level=1)
                for line in lines[1:]:
                    doc.add_paragraph(line, style='List Bullet')
        doc.save(docx_path)

        # Generate slide-specific outline
        slide_prompt = f"""
Create a training presentation titled “{prompt}” with 5 to 7 clearly defined slides.

Each slide should include:
- A clear slide title
- 3 to 5 bullet points
- Content focused on educating a beginner-to-intermediate audience

Use simple language, short phrases, and group related ideas together.

Format the output as:

Slide 1: Title
- Bullet 1
- Bullet 2

Slide 2: Title
- Bullet 1
- Bullet 2

End the outline when complete.
"""

        try:
            slide_response = llm.invoke(slide_prompt)
            slide_outline = slide_response.content if hasattr(slide_response, "content") else str(slide_response)
        except Exception:
            slide_outline = llm.call_as_llm(slide_prompt)

        slide_outline = slide_outline.strip()

        # Save .pptx
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0]).shapes.title.text = prompt

        for section in slide_outline.strip().split("Slide")[1:]:
            lines = section.strip().split("\n")
            if len(lines) > 1:
                title_line = lines[0].replace(":", "").strip()
                bullet_points = [line.strip("- ").strip() for line in lines[1:] if line.strip()]
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = title_line
                textbox = slide.placeholders[1]
                textbox.text = "\n".join(bullet_points)

        prs.save(pptx_path)

        return render_template(
            "course_result.html",
            result=course_outline,
            prompt=prompt,
            category=course_category,
            course_file=txt_file,
            docx_file=docx_file,
            pptx_file=pptx_file,
            pdf_file=pdf_file
        )

    return render_template("generate_course.html")
