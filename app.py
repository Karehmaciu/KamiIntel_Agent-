# Standard library imports
import os
import re
import csv
import logging
from datetime import datetime

# Third-party imports
from flask import (
    Flask, request, render_template, send_file, redirect, 
    url_for, flash, send_from_directory, abort
)
from werkzeug.utils import secure_filename
from langchain_openai import ChatOpenAI, OpenAIEmbeddings
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
import pandas as pd
import magic  # python-magic-bin on Windows

# Local application imports
from course_generator import generate_course_outline
from train_handler import train_on_any_file, VECTOR_DB_PATH
from agent_utils.memory_manager import (
    save_prompt_and_response, export_to_word,
    generate_pptx_advanced, generate_pptx_slides_from_text
)
from agent_utils.vector_store import query_vector_store, train_on_csv
from file_security import is_file_safe, is_file_size_safe, sanitize_filename

app = Flask(__name__)
app.secret_key = "your-secret-key"  # Consider using environment variables for secrets

# Setup basic logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    filename=os.path.join('chat_data', 'logs', 'application.log'),
    filemode='a'
)
logger = logging.getLogger('kamiintel')

GENERATED_REPORT_DIR = "chat_data/generated_reports"
os.makedirs(GENERATED_REPORT_DIR, exist_ok=True)

@app.route("/generate_course", methods=["GET", "POST"])
def generate_course():
    if request.method == "POST":
        prompt = request.form.get("course_prompt", "")
        category = request.form.get("category")
        
        logger.info(f"Course generation requested for topic: '{prompt}', category: '{category}'")

        if not prompt:
            logger.warning("Empty course prompt submitted")
            flash("⚠️ Please enter a course topic.", "danger")
            return redirect(request.url)

        # Generate course content
        try:
            result = generate_course_outline(prompt)
            if isinstance(result, str):
                logger.error(f"Course generation error: {result}")
                flash(result, "danger")
                return redirect(request.url)
            
            logger.info(f"Course generated successfully for: '{prompt}'")
            
            course_outline = result["content"]
            course_category = category or result["category"]
            date_str = datetime.today().strftime('%Y-%m-%d')
            safe_topic = prompt.lower().replace(" ", "_").replace("/", "-")
            base_filename = f"{safe_topic}_{date_str}"

            # Define paths
            save_dir = "chat_data/generated_courses"
            os.makedirs(save_dir, exist_ok=True)
            txt_file = f"{base_filename}.txt"
            docx_file = f"{base_filename}.docx"
            pptx_file = f"{base_filename}.pptx"
            pdf_file = f"{base_filename}.pdf"  # placeholder

            txt_path = os.path.join(save_dir, txt_file)
            docx_path = os.path.join(save_dir, docx_file)
            pptx_path = os.path.join(save_dir, pptx_file)

            # Save TXT
            with open(txt_path, "w", encoding="utf-8") as f:
                f.write(f"{prompt}\nCategory: {course_category}\nGenerated on: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n{course_outline}")

            # Save DOCX
            doc = Document()
            doc.add_heading(prompt, 0)
            doc.add_paragraph(f"Category: {course_category}")
            doc.add_paragraph(f"Generated on: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}") 
            doc.add_paragraph("")
            for section in course_outline.strip().split("\n\n"):
                lines = section.strip().split("\n")
                if lines:
                    doc.add_heading(lines[0], level=1)
                    for line in lines[1:]:
                        doc.add_paragraph(line, style='List Bullet')
            doc.save(docx_path)

            # Save PPTX
            prs = Presentation()
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = prompt
            slide.placeholders[1].text = f"Category: {course_category}\nGenerated on {datetime.now().strftime('%Y-%m-%d')}"

            for section in course_outline.strip().split("\n\n"):
                lines = section.strip().split("\n")
                if lines:
                    slide = prs.slides.add_slide(slide_layout)
                    slide.shapes.title.text = lines[0]
                    content = "\n".join(lines[1:])
                    slide.placeholders[1].text = content

            prs.save(pptx_path)

            return render_template(
                "course_result.html",
                result=course_outline,
                prompt=prompt,
                course_file=txt_file,
                docx_file=docx_file,
                pptx_file=pptx_file,
                pdf_file=pdf_file,
                category=course_category
            )
        except Exception as e:
            logger.exception(f"Unexpected error during course generation: {e}")
            flash("An error occurred while generating the course. Please try again.", "danger")
            return redirect(request.url)

    return render_template("generate_course.html")

@app.route("/courses")
def view_courses():
    course_folder = "chat_data/generated_courses"
    courses = []
    
    logger.info("Viewing course library")

    if os.path.exists(course_folder):
        for root, dirs, files in os.walk(course_folder):
            for fname in sorted(files):
                # Skip system files and hidden files
                if fname.startswith('.') or fname.startswith('~$'):
                    continue
                    
                full_path = os.path.join(root, fname)
                rel_path = os.path.relpath(full_path, start="chat_data")
                time_str = datetime.fromtimestamp(os.path.getctime(full_path)).strftime("%Y-%m-%d %H:%M")
                
                # Get file type for better display
                file_type = os.path.splitext(fname)[1].lstrip('.').upper()
                
                # Get category from folder structure
                category_path = os.path.relpath(root, course_folder)
                category = category_path if category_path != '.' else 'General'
                
                courses.append({
                    "name": fname,
                    "time": time_str,
                    "url": "/chat_data/" + rel_path.replace("\\", "/"),
                    "type": file_type,
                    "category": category.replace('_', ' ').title()
                })
    
    return render_template("view_courses.html", courses=courses)

@app.route("/view_course/<category>/<filename>")
def view_course_detail(category, filename):
    filepath = os.path.join("chat_data", "generated_courses", category, filename)
    if not os.path.exists(filepath):
        return "Course not found", 404

    with open(filepath, "r", encoding="utf-8") as f:
        content = f.read()

    return render_template("course_detail.html", content=content, filename=filename, category=category)



@app.route("/delete_course/<filename>")
def delete_course(filename):
    # Basic security check - prevent path traversal
    if '..' in filename or '/' in filename:
        logger.warning(f"Potential path traversal attempt: {filename}")
        flash("Invalid filename", "danger")
        return redirect("/courses")
        
    course_path = os.path.join("chat_data/generated_courses", filename)
    if os.path.exists(course_path):
        try:
            os.remove(course_path)
            logger.info(f"Deleted course file: {filename}")
            flash(f"File {filename} deleted successfully", "success")
        except Exception as e:
            logger.error(f"Error deleting course file {filename}: {e}")
            flash(f"Error deleting file: {e}", "danger")
    else:
        flash("File not found", "warning")
        
    return redirect("/courses")

@app.route('/slides/<filename>')
def download_slide(filename):
    return send_from_directory('static/slides', filename)


@app.route("/train", methods=["GET", "POST"])
def upload_and_train_file():
    uploaded_files = []

    if request.method == "POST":
        file = request.files.get("ffd_file")
        if not file or file.filename == "":
            flash("⚠️ No file selected.", "danger")
            return redirect(request.url)

        # 🧼 Sanitize filename
        safe_filename = sanitize_filename(file.filename)
        upload_dir = "chat_data/uploads"
        os.makedirs(upload_dir, exist_ok=True)
        save_path = os.path.join(upload_dir, safe_filename)
        file.save(save_path)

        # 🛡️ Check MIME type
        if not is_file_safe(save_path):
            os.remove(save_path)
            flash("❌ Rejected: Suspicious file type.", "danger")
            return redirect(request.url)

        # 🧱 Check file size
        if not is_file_size_safe(save_path):
            os.remove(save_path)
            flash("❌ Rejected: File is too large (max 10MB).", "danger")
            return redirect(request.url)

        # ✅ Proceed with training
        try:
            result = train_on_any_file(save_path)
            flash(f"✅ File '{safe_filename}' uploaded and trained successfully!", "success")
        except Exception as e:
            flash(f"❌ Error training on file '{safe_filename}': {e}", "danger")

        return redirect(request.url)

    if os.path.exists("chat_data/uploads"):
        uploaded_files = os.listdir("chat_data/uploads")

    return render_template("train_ffd.html", uploaded_files=uploaded_files)




@app.route("/", methods=["GET", "POST"])
def chat():
    response = ""
    prompt = ""
    language = "en"
    if request.method == "POST":
        if "user_prompt" not in request.form:
            return render_template("chat.html", response="⚠️ No prompt received.", prompt="")

        prompt = request.form["user_prompt"]
        language = request.form.get("language", "en")
        action = request.form.get("action")
        lang_suffix = f"_{language}" if language != "en" else ""

        if "open excel and write a report about" in prompt.lower():
            # CSV report logic...
            topic = prompt.lower().split("write a report about", 1)[-1].strip()
            llm = ChatOpenAI(temperature=0.3, model_name="gpt-3.5-turbo")
            formatted_prompt = f"List counties in Kenya and their populations in a clean CSV format. Use only County and Population as headers. Do not include pipes or Markdown formatting. Respond in {language}."

            try:
                data = llm.call_as_llm(formatted_prompt) 
                #data = llm.invoke(formatted_prompt)

                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                filename = os.path.join(GENERATED_REPORT_DIR, f"csv_report_{timestamp}{lang_suffix}.csv")

                with open(filename, mode="w", newline="", encoding="utf-8") as csvfile:
                    writer = csv.writer(csvfile)
                    for i, line in enumerate(data.strip().split("\n")):
                        clean_line = line.strip(" |").replace(",", "")
                        if i == 0:
                            if "CountyPopulation" in clean_line:
                                writer.writerow(["County", "Population"])
                            else:
                                writer.writerow([col.strip() for col in clean_line.split("|") if col])
                        else:
                            match = re.match(r"([A-Za-z\s]+)(\d+)", clean_line)
                            if match:
                                writer.writerow([match.group(1).strip(), match.group(2)])

                flash("CSV Report generated successfully!", "success")
                return redirect(url_for("chat", download=filename.replace("\\", "/")))
            except Exception as e:
                flash(f"❌ Error generating CSV report: {e}", "danger")

        elif action == "generate_pptx" or "generate powerpoint slides about" in prompt.lower():
            topic = prompt.lower().split("about", 1)[-1].strip() if "about" in prompt.lower() else prompt.strip()
            llm = ChatOpenAI(temperature=0.4, model_name="gpt-3.5-turbo")
            formatted_prompt = f"Create a slide outline for a presentation about {topic}. List 5 slide titles with 2-3 bullet points each. Use {language}."

            try:
                content = llm.call_as_llm(formatted_prompt)
                #data = llm.invoke(formatted_prompt)

                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                pptx_filename = os.path.join(GENERATED_REPORT_DIR, f"slides_{timestamp}{lang_suffix}.pptx")
                generate_pptx_slides_from_text(content, topic, pptx_filename)

                flash("PowerPoint slides successfully generated!", "success")
                return redirect(url_for("chat", download=pptx_filename.replace("\\", "/")))
            except Exception as e:
                flash(f"❌ Error generating slides: {e}", "danger")

        elif "upload csv and summarize" in prompt.lower():
            try:
                upload_path = "chat_data/raw_data/latest_upload.csv"
                try:
                    df = pd.read_csv(upload_path, encoding="utf-8")
                except UnicodeDecodeError:
                    df = pd.read_csv(upload_path, encoding="ISO-8859-1")
                train_on_csv(upload_path)
                summary = df.describe(include="all").to_string()
                response = f"📊 Summary of uploaded CSV:\n\n{summary}"
            except Exception as e:
                flash(f"❌ Error analyzing CSV upload: {e}", "danger")

        elif any(x in prompt.lower() for x in ["write a report about", "andika ripoti", "écrire un rapport", "berichten", "የ", "ripoti"]):
            topic = prompt.lower().split("about")[-1].strip() if "about" in prompt.lower() else prompt.strip()
            llm = ChatOpenAI(temperature=0.4, model_name="gpt-3.5-turbo")
            formatted_prompt = f"Write a professional report about: {topic}. Include title, introduction, 3 key sections, and a conclusion. Use {language}."

            try:
                report = llm.call_as_llm(formatted_prompt)
                #data = llm.invoke(formatted_prompt)

                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                filename = os.path.join(GENERATED_REPORT_DIR, f"report_{timestamp}{lang_suffix}.docx")

                doc = Document()
                doc.add_heading(f"{topic.title()} Report", 0)
                for line in report.split('\n'):
                    if line.strip():
                        doc.add_paragraph(line.strip())
                doc.save(filename)

                flash("Word report successfully generated!", "success")
                return redirect(url_for("chat", download=filename.replace("\\", "/")))

            except Exception as e:
                flash(f"❌ Error generating Word report: {e}", "danger")

        elif action == "regenerate":
            response = query_vector_store(prompt, use_chat_history=False)
            if not response or "No relevant" in response:
                llm = ChatOpenAI(temperature=0.5, model_name="gpt-3.5-turbo")
                retry_response = llm.invoke(prompt)
                response = retry_response.strip() if retry_response.strip() else "👶 ELI5: We didn’t find useful documents for that. Try rephrasing or uploading more content."
        
        
        elif action == "explain_like_5":
            context = query_vector_store(prompt, use_chat_history=False)
            if not context or "No relevant" in context:
                 response = "👶 ELI5: We couldn’t find enough information to explain. Try uploading more data or asking something else."
            else:
                llm = ChatOpenAI(temperature=0.4, model_name="gpt-3.5-turbo")
                simplified = llm.call_as_llm(f"Explain this like I'm 5 years old:\n{context}")
                response = f"👶 ELI5: {simplified.content.strip()}"

        else:
            response = query_vector_store(prompt)
            if not response or "No relevant" in response:
                llm = ChatOpenAI(temperature=0.5, model_name="gpt-3.5-turbo")
                retry_response = llm.invoke(prompt)
                response = retry_response.content.strip() if retry_response.content.strip() else "👶 ELI5: We didn’t find useful documents for that. Try rephrasing or uploading more content."
            
            save_prompt_and_response(prompt, response)

                 # Automatically save a copy of the response to generated_reports/
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            response_filename = os.path.join(GENERATED_REPORT_DIR, f"response_{timestamp}.docx")


            doc = Document()
            title = prompt.strip().split("\n")[0][:40]  # First line or first 40 characters
            doc.add_heading(f"Response: {title}", 0)
            doc.add_paragraph("Prompt:")
            doc.add_paragraph(prompt)
            doc.add_paragraph("\nResponse:")
            doc.add_paragraph(response)
            doc.save(response_filename)


        if "export" in request.form:
            filename = export_to_word(prompt, response)
            return send_file(filename, as_attachment=True)

    download_file = request.args.get("download")
    return render_template("chat.html", response=response, prompt=prompt, download_file=download_file, language=language)

@app.route("/slides", methods=["GET", "POST"])
def slides():
    download_file = None

    if request.method == "POST":
        topic = request.form.get("topic", "Untitled")
        language = request.form.get("language", "en")
        theme_color = request.form.get("theme_color", "#ffffff").lstrip("#")
        text_color = request.form.get("text_color", "#282850").lstrip("#")
        image_file = request.files.get("image_file")
        image_url = request.form.get("image_url", "")
        chart_data = request.form.get("chart_data", "")
        max_slides = int(request.form.get("max_slides", 5))

        image_path = None
        if image_file and image_file.filename:
            image_path = os.path.join("chat_data", "raw_data", image_file.filename)
            image_file.save(image_path)

        llm = ChatOpenAI(temperature=0.4, model_name="gpt-3.5-turbo")
        formatted_prompt = f"Create a slide outline for a presentation about {topic}. List 10 slide titles with 2-3 bullet points each. Use {language}."
        #content = llm.invoke(formatted_prompt)
        response = llm.invoke(formatted_prompt)
        content = response.content if hasattr(response, "content") else response

        # Build slide deck
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"slides_advanced_{timestamp}_{language}.pptx"
        filepath = os.path.join(GENERATED_REPORT_DIR, filename) 

        theme_rgb = tuple(int(theme_color[i:i+2], 16) for i in (0, 2, 4))
        text_rgb = tuple(int(text_color[i:i+2], 16) for i in (0, 2, 4))

        generate_pptx_advanced(
            content=content,
            topic=topic,
            save_path=filepath,
            image_path=image_path,
            image_url=image_url,
            chart_input=chart_data,
            theme_color=theme_rgb,
            text_color=text_rgb,
            max_slides=max_slides,
            language=language
        )

        flash("✅ Slides generated successfully!", "success")
        download_file = filename

    return render_template("slides.html", download_file=download_file)


@app.route("/upload", methods=["GET", "POST"])
def upload():
    if request.method == "POST":
        if "file" not in request.files:
            flash("No file part", "danger") 
            return redirect(request.url)

        file = request.files.get("file")
        if not file or file.filename == "":
            flash("No selected file", "danger")
            return redirect(request.url)

        # Process the uploaded file
        upload_dir = "chat_data/raw_data"
        os.makedirs(upload_dir, exist_ok=True)
        file_path = os.path.join(upload_dir, "latest_upload.csv")
        file.save(file_path)

        try:
            train_on_csv(file_path)
            flash("✅ File uploaded and used for training!", "success")
        except Exception as e:
            flash(f"⚠️ File uploaded but failed to train agent: {e}", "danger")

        return redirect(url_for("chat"))

    return render_template("upload.html")

@app.route("/reports")
def view_reports():
    reports = []
    if os.path.exists(GENERATED_REPORT_DIR):
        for fname in os.listdir(GENERATED_REPORT_DIR):
            if fname.endswith((".docx", ".csv", ".pptx")):
                fpath = os.path.join(GENERATED_REPORT_DIR, fname)
                ctime = os.path.getctime(fpath)
                language = "en"
                if "_sw" in fname:
                    language = "sw"
                elif "_fr" in fname:
                    language = "fr"
                elif "_de" in fname:
                    language = "de"
                elif "_am" in fname:
                    language = "am"

                title = os.path.splitext(fname)[0].replace("_", " ").title()
                
                # Don't try to read binary files as text
                content_preview = ""
                if fname.endswith(".csv"):
                    try:
                        with open(fpath, "r", encoding="utf-8", errors="ignore") as f:
                            content_preview = f.read(1000)
                    except Exception:
                        content_preview = "Preview not available"

                reports.append({
                    "name": fname,
                    "title": title,
                    "language": language,
                    "path": url_for("download_file", filename=fname),
                    "time": datetime.fromtimestamp(ctime).strftime("%Y-%m-%d %H:%M:%S"),
                    "preview": content_preview,
                    "type": fname.split(".")[-1].upper()
                })

        reports.sort(key=lambda x: x['time'], reverse=True)

    return render_template("reports.html", reports=reports) 

@app.route("/train_all_uploads")
def train_all_uploads():
    upload_dir = "chat_data/uploads"
    trained_files = []
    skipped_files = []

    if not os.path.exists(upload_dir):
        flash("⚠️ No uploaded files found to train.", "danger")
        return redirect(url_for("upload_and_train_file"))

    files = os.listdir(upload_dir)
    if not files:
        flash("⚠️ Upload folder is empty.", "warning")
        return redirect(url_for("upload_and_train_file"))

    for filename in files:
        file_path = os.path.join(upload_dir, filename)

        if not is_file_safe(file_path) or not is_file_size_safe(file_path):
            skipped_files.append(filename)
            continue

        try:
            train_on_any_file(file_path)
            trained_files.append(filename)
        except Exception as e:
            skipped_files.append(filename)
            print(f"❌ Error training on {filename}: {e}")

    if trained_files:
        flash(f"✅ Trained on {len(trained_files)} file(s): {', '.join(trained_files)}", "success")
    if skipped_files:
        flash(f"⚠️ Skipped {len(skipped_files)} file(s) due to issues: {', '.join(skipped_files)}", "warning")

    return redirect(url_for("upload_and_train_file"))
 

@app.route("/download/<filename>")
def download_file(filename):
    file_path = os.path.join(GENERATED_REPORT_DIR, filename)
    if os.path.exists(file_path):
        # Set appropriate MIME types for docx and txt files
        mimetype = None
        if filename.endswith('.docx'):
            mimetype = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        elif filename.endswith('.txt'):
            mimetype = 'text/plain'
        return send_file(file_path, as_attachment=True, mimetype=mimetype)
    return "File not found.", 404

@app.route("/download_course/<category>/<filename>")
def download_course(category, filename):
    file_path = os.path.join("chat_data/generated_courses", category, filename)
    if os.path.exists(file_path):
        # Set appropriate MIME types for docx and txt files
        mimetype = None
        if filename.endswith('.docx'):
            mimetype = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        elif filename.endswith('.txt'):
            mimetype = 'text/plain'
        return send_file(file_path, as_attachment=True, mimetype=mimetype)
    return "File not found", 404

@app.route("/chat_data/<path:filename>") 
def serve_chatdata_file(filename):
    """Serve files from the chat_data directory."""
    # Create a safe path by ensuring we only access files within chat_data
    file_path = os.path.join("chat_data", filename)
    if os.path.exists(file_path) and os.path.isfile(file_path):
        # Set appropriate MIME types for docx and txt files
        mimetype = None
        if filename.endswith('.docx'):
            mimetype = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        elif filename.endswith('.txt'):
            mimetype = 'text/plain'
        return send_file(file_path, as_attachment=True, mimetype=mimetype)
    return "File not found.", 404


@app.route("/course_portal") 
def course_portal():
    """Course Generator Portal page with expanded course generation features"""
    return render_template("generate_course.html")

if __name__ == "__main__":
    os.makedirs("chat_data/prompts", exist_ok=True)
    os.makedirs("chat_data/responses", exist_ok=True)
    os.makedirs("chat_data/raw_data", exist_ok=True)
    os.makedirs("chat_data/logs", exist_ok=True)
    os.makedirs(GENERATED_REPORT_DIR, exist_ok=True)

# ✅ Only one main entry point
if __name__ == "__main__":
    print("✅ About to start server on port 8002...")
    import socket
    print(f"🧪 Test - Local IP: http://{socket.gethostbyname(socket.gethostname())}:8002")    
    app.run(debug=True, port=8002, host="0.0.0.0") 
