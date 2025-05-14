import os
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import pandas as pd
from langchain_community.embeddings import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.schema import Document as LCDocument

VECTOR_DB_PATH = "chat_data/vector_store"

def extract_text_from_file(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    if ext == ".txt":
        with open(filepath, "r", encoding="utf-8") as f:
            return f.read()
    elif ext == ".docx":
        doc = Document(filepath)
        return "\n".join([p.text for p in doc.paragraphs])
    elif ext == ".pdf":
        doc = fitz.open(filepath)
        return "\n".join([page.get_text() for page in doc])
    elif ext == ".pptx":
        prs = Presentation(filepath)
        return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    elif ext in [".xlsx", ".xls"]:
        df = pd.read_excel(filepath, sheet_name=None)
        return "\n".join([str(cell) for sheet in df.values() for row in sheet.values for cell in row if pd.notna(cell)])
    else:
        return "Unsupported file format."

def train_on_any_file(filepath):
    text = extract_text_from_file(filepath)
    if text.startswith("Unsupported"):
        return text
    docs = [LCDocument(page_content=text)]
    embeddings = OpenAIEmbeddings()
    db = FAISS.from_documents(docs, embeddings)
    db.save_local(VECTOR_DB_PATH)
    return "✅ File trained and stored in vector DB."
